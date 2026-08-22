import sys
from pptx import Presentation
from pptx.util import Pt

def replace_text(shape, new_text):
    if shape.has_text_frame:
        # Keep the original formatting of the first run if possible
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            font_name = shape.text_frame.paragraphs[0].runs[0].font.name
            font_size = shape.text_frame.paragraphs[0].runs[0].font.size
            font_bold = shape.text_frame.paragraphs[0].runs[0].font.bold
            font_color = shape.text_frame.paragraphs[0].runs[0].font.color.rgb if shape.text_frame.paragraphs[0].runs[0].font.color.type == 1 else None
        else:
            font_name, font_size, font_bold, font_color = None, None, None, None

        shape.text = new_text
        
        # Apply formatting back
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if font_name: r.font.name = font_name
                if font_size: r.font.size = font_size
                if font_bold: r.font.bold = font_bold
                if font_color: r.font.color.rgb = font_color

def populate_ppt(template_path, output_path):
    prs = Presentation(template_path)
    
    # Slide 1: INCUBATEE DETAILS WITH TEAM MEMBERS
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame and shape != s1.shapes.title:
            replace_text(shape, "Team Lead: SANTHIVARSHINI D\n\nTEAM MEMBERS: [Please fill your team members]\n\nMENTOR DETAILS: [Please fill your mentor details]")

    # Slide 2: THEME
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame and shape != s2.shapes.title:
            replace_text(shape, "Theme/Sector: Miscellaneous Sector (IT, Automation, Machinery)\n\nTitle: FactoryDNA - AI-Powered Manufacturing Knowledge Intelligence Platform for MSMEs")

    # Slide 3: CONCEPT & OBJECTIVE OF THE IDEA
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.has_text_frame and shape != s3.shapes.title:
            replace_text(shape, "FactoryDNA is an AI-powered Manufacturing Knowledge Intelligence Platform designed to help MSMEs preserve, organize and utilize valuable operational knowledge often lost due to employee turnover.\n\nIt continuously collects information from production records, maintenance logs, sensor data, SOPs, and voice interactions to build a centralized knowledge repository.\n\nObjective: Reduce knowledge loss, improve decision-making, and accelerate workforce training via an intelligent Manufacturing Copilot. It enhances production efficiency, product quality, and supports predictive maintenance.")

    # Slide 5: Working principle
    s4 = prs.slides[4]
    for shape in s4.shapes:
        if shape.has_text_frame and shape != s4.shapes.title:
            replace_text(shape, "1. Data Collection: Captures unstructured info (machine logs, maintenance records, SOPs, images, voice interactions).\n\n2. AI Integration: Utilizes Artificial Intelligence, Knowledge Graphs, Retrieval-Augmented Generation (RAG), Edge AI, Computer Vision, and IoT.\n\n3. Knowledge Base: Converts the collected data into a structured, searchable manufacturing knowledge base.\n\n4. AI Copilot: Provides contextual guidance, troubleshooting recommendations, and process optimization based on the factory's own historical knowledge.")

    # Slide 6: Noval & Innovative Value
    s5 = prs.slides[5]
    for shape in s5.shapes:
        if shape.has_text_frame and shape != s5.shapes.title:
            replace_text(shape, "Unlike existing Industry 4.0 solutions that only monitor machines or production data, FactoryDNA focuses on capturing human expertise—the experience, decision-making patterns, and best practices of skilled operators and engineers.\n\nBy transforming human expertise into reusable digital intelligence, it ensures that when experienced employees retire, their knowledge remains accessible to the organization.")

    # Slide 7: Innovative Solutions for a Better and Sustainable Society
    s6 = prs.slides[6]
    for shape in s6.shapes:
        if shape.has_text_frame and shape != s6.shapes.title:
            replace_text(shape, "- Reduces dependency on individual experts and shortens training time for new workforce.\n- Improves product quality and preserves organizational knowledge.\n- Affordable modular deployment supports productivity, compliance, and sustainability in small and medium manufacturing units.\n- Ensures long-term competitiveness by building resilient operations.")

    # Slide 8: Innovative value proposition of products/ services being offered
    s7 = prs.slides[7]
    for shape in s7.shapes:
        if shape.has_text_frame and shape != s7.shapes.title:
            replace_text(shape, "- AI-driven insights: An intelligent Manufacturing Copilot answers process queries and recommends best practices.\n- Continuous learning: The platform learns from production outcomes to update recommended practices.\n- Operational continuity: Makes critical expertise available across shifts and generations of workers.")

    # Slide 9: Ease of availability of the resources used in the product
    s8 = prs.slides[8]
    for shape in s8.shapes:
        if shape.has_text_frame and shape != s8.shapes.title:
            replace_text(shape, "- Scalable Software-as-a-Service (SaaS) model with optional Edge AI deployment.\n- Integrates seamlessly with existing ERP, MES, IoT sensors, and production systems without requiring major infrastructure changes.\n- Modular architecture allows MSMEs to adopt features based on their specific business requirements.")

    # Slide 10: Strategic Benefits
    s9 = prs.slides[9]
    for shape in s9.shapes:
        if shape.has_text_frame and shape != s9.shapes.title:
            replace_text(shape, "Commercial/Financial Viability:\n- Strong revenue opportunities via subscription plans, enterprise licensing, AI Copilot services, and analytics modules.\n\nStrategic Alignment:\n- Aligns with Digital India, Make in India, Industry 4.0, Industry 5.0, Skill India, and Atmanirbhar Bharat.\n\nScalability:\n- Significant potential in developing economies (Asia, Africa, Latin America) facing similar workforce knowledge loss challenges.")

    # Slide 11: Impact on the project
    s10 = prs.slides[10]
    for shape in s10.shapes:
        if shape.has_text_frame and shape != s10.shapes.title:
            replace_text(shape, "Market Impact:\n- Addresses a critical gap for India's 6 crore+ MSMEs that depend on undocumented expertise.\n- Broad applicability across diverse sectors (auto, CNC, textiles, pharma, FMCG, etc.).\n\nOperational Impact:\n- Improves efficiency, strengthens compliance, and reduces knowledge transfer friction.")

    prs.save(output_path)
    print(f"Presentation successfully generated and saved to {output_path}")

if __name__ == "__main__":
    populate_ppt(sys.argv[1], sys.argv[2])
