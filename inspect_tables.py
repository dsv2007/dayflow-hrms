import sys
from pptx import Presentation

def inspect_shapes(slide):
    for i, shape in enumerate(slide.shapes):
        print(f"Shape {i}: Type: {shape.shape_type}")
        if shape.has_text_frame:
            print(f"Text: {shape.text[:50]}")
        if shape.has_table:
            print(f"Table: {len(shape.table.rows)}x{len(shape.table.columns)}")

if __name__ == "__main__":
    prs = Presentation(sys.argv[1])
    print("--- Slide 1 ---")
    inspect_shapes(prs.slides[0])
    print("--- Slide 2 ---")
    inspect_shapes(prs.slides[1])
    print("--- Slide 7 ---")
    inspect_shapes(prs.slides[6])
    print("--- Slide 10 ---")
    inspect_shapes(prs.slides[9])
    print("--- Slide 11 ---")
    inspect_shapes(prs.slides[10])
    print("--- Slide 12 ---")
    inspect_shapes(prs.slides[11])
