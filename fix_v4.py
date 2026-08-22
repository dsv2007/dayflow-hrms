import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def center_text(shape):
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

def create_block(slide, left, top, width, height, text, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, rgb=(0, 102, 204)):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(1.5)
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_down_arrow(slide, left, top, width=Inches(0.4), height=Inches(0.4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shape.line.color.rgb = RGBColor(100, 100, 100)

def fix_presentation(template_path, output_path):
    prs = Presentation(template_path)
    
    # ---------------------------------------------
    # Fix Flow Chart on Slide 4 (index 3)
    # ---------------------------------------------
    slide4 = prs.slides[3]
    
    # DELETE ALL SHAPES EXCEPT TITLE
    shapes_to_delete = []
    for shape in slide4.shapes:
        if shape != slide4.shapes.title:
            shapes_to_delete.append(shape)
            
    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)
        
    # Create new aligned flowchart
    width = Inches(6.0)
    left = Inches(2.0)
    height = Inches(0.9)
    
    # 1. Inputs
    top1 = Inches(1.5)
    create_block(slide4, left, top1, width, height, 
                 "1. Data Inputs\nUnstructured (Logs, SOPs) & Structured (IoT, ERP)", 
                 rgb=(39, 174, 96)) # Green
    
    # Arrow 1
    add_down_arrow(slide4, Inches(4.8), Inches(2.45), Inches(0.4), Inches(0.5))
    
    # 2. AI Engine
    top2 = Inches(3.0)
    create_block(slide4, left, top2, width, height, 
                 "2. AI Processing Engine\n(RAG, Edge AI, Knowledge Graph)", 
                 rgb=(211, 84, 0)) # Orange
                 
    # Arrow 2
    add_down_arrow(slide4, Inches(4.8), Inches(3.95), Inches(0.4), Inches(0.5))
    
    # 3. Knowledge Base
    top3 = Inches(4.5)
    create_block(slide4, left, top3, width, height, 
                 "3. Centralized Knowledge Base\n(Searchable & Structured Repository)", 
                 rgb=(41, 128, 185)) # Blue
                 
    # Arrow 3
    add_down_arrow(slide4, Inches(4.8), Inches(5.45), Inches(0.4), Inches(0.5))
    
    # 4. Copilot
    top4 = Inches(6.0)
    create_block(slide4, left, top4, width, height, 
                 "4. AI Manufacturing Copilot\n(Contextual Guidance & Optimization)", 
                 shape_type=MSO_SHAPE.HEXAGON, rgb=(142, 68, 173)) # Purple

    # ---------------------------------------------
    # Fix Budget Table on Slide 12 (index 11)
    # ---------------------------------------------
    slide12 = prs.slides[11]
    for shape in slide12.shapes:
        if shape.has_table:
            table = shape.table
            
            # Helper to set text and font size
            def set_cell_text(row, col, text):
                cell = table.cell(row, col)
                cell.text = text
                # Force font size to match template (usually around 14pt)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for r in p.runs:
                        r.font.size = Pt(14)
            
            # Machine Usage
            set_cell_text(1, 3, "2.0")
            # Electricity
            set_cell_text(2, 3, "0.5")
            # Raw material (IoT, sensors)
            set_cell_text(3, 3, "3.5")
            # Testing/calibration
            set_cell_text(4, 3, "1.0")
            # Other charges
            set_cell_text(5, 3, "3.0")
            
            # Mentor
            set_cell_text(6, 3, "1.0")
            
            # Travel
            set_cell_text(7, 3, "1.0")
            
    prs.save(output_path)
    print(f"Presentation successfully updated and saved to {output_path}")

if __name__ == "__main__":
    fix_presentation(sys.argv[1], sys.argv[2])
