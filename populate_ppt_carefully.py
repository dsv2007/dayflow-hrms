import sys
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def get_body_shape(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            return shape
    return None

def append_to_paragraph(paragraph, text, bold=False):
    run = paragraph.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    # try to match size
    if len(paragraph.runs) > 1 and paragraph.runs[0].font.size:
        run.font.size = paragraph.runs[0].font.size

def populate_ppt_carefully(template_path, output_path):
    prs = Presentation(template_path)
    
    # Slide 1: INCUBATEE DETAILS
    s1 = prs.slides[0]
    body1 = get_body_shape(s1)
    if body1:
        for p in body1.text_frame.paragraphs:
            if "TEAM LEAD:" in p.text.upper():
                append_to_paragraph(p, " SANTHIVARSHINI D", bold=True)
            elif "TEAM MEMBERS:" in p.text.upper():
                append_to_paragraph(p, " [Insert Team Members]", bold=True)
            elif "MENTOR DETAILS:" in p.text.upper():
                append_to_paragraph(p, " [Insert Mentor Details]", bold=True)
                
    # Slide 2: THEME Table
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_table:
            table = shape.table
            # Row 6 is "Other Frontier Technologies"
            cell = table.cell(6, 1)
            p = cell.text_frame.paragraphs[0]
            append_to_paragraph(p, " -> FactoryDNA (Miscellaneous: IT, Automation)", bold=True)
            # highlight background
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 0) # yellow highlight

    # Slide 3: CONCEPT & OBJECTIVE
    s3 = prs.slides[2]
    body3 = get_body_shape(s3)
    if body3:
        body3.text = "FactoryDNA is an AI-powered Manufacturing Knowledge Intelligence Platform designed to help MSMEs preserve, organize and utilize valuable operational knowledge often lost due to employee turnover.\n\nObjective: Reduce knowledge loss, improve decision-making, and accelerate workforce training via an intelligent Manufacturing Copilot. It enhances production efficiency, product quality, and supports predictive maintenance."

    # Slide 5: Working principle
    s5 = prs.slides[4]
    body5 = get_body_shape(s5)
    if body5:
        body5.text = "1. Data Collection: Captures unstructured info (machine logs, maintenance records, SOPs, images, voice interactions).\n\n2. AI Integration: Utilizes Artificial Intelligence, Knowledge Graphs, RAG, Edge AI, and Computer Vision.\n\n3. AI Copilot: Provides contextual guidance, troubleshooting recommendations, and process optimization based on historical knowledge."

    # Slide 6: Noval & Innovative Value
    s6 = prs.slides[5]
    body6 = get_body_shape(s6)
    if body6:
        body6.text = "Unlike existing Industry 4.0 solutions that only monitor machines, FactoryDNA focuses on capturing human expertise—the experience, decision-making patterns, and best practices of skilled operators.\n\nIt transforms human expertise into reusable digital intelligence."

    # Slide 7: Innovative Solutions
    s7 = prs.slides[6]
    body7 = get_body_shape(s7)
    if body7:
        for p in body7.text_frame.paragraphs:
            if "Useful to Society" in p.text:
                p.text = "Useful to Society: Preserves organizational knowledge and shortens training time for the new workforce."
            elif "Reduce cost of living" in p.text:
                p.text = "Reduce cost of production: Affordable deployment for MSMEs, improving product quality and reducing dependency on individual experts."
            elif "Improve quality of life" in p.text:
                p.text = "Improve work quality: Empowers workers with an AI Copilot for continuous learning and standardized operations."
            elif "energy Initiative" in p.text:
                p.text = "Sustainability: Supports long-term competitiveness and resilient operations without major infrastructure changes."

    # Slide 8: Innovative value proposition
    s8 = prs.slides[7]
    body8 = get_body_shape(s8)
    if body8:
        body8.text = "AI-driven insights: An intelligent Manufacturing Copilot answers process queries and recommends best practices.\n\nContinuous learning: The platform learns from production outcomes to update recommended practices.\n\nOperational continuity: Makes critical expertise available across shifts and generations of workers."

    # Slide 9: Ease of availability of resources
    s9 = prs.slides[8]
    body9 = get_body_shape(s9)
    if body9:
        body9.text = "Scalable Software-as-a-Service (SaaS) model with optional Edge AI deployment.\n\nIntegrates seamlessly with existing ERP, MES, IoT sensors, and production systems without requiring major infrastructure changes.\n\nModular architecture allows MSMEs to adopt features based on their specific business requirements."

    # Slide 10: Strategic Benefits
    s10 = prs.slides[9]
    body10 = get_body_shape(s10)
    if body10:
        for p in body10.text_frame.paragraphs:
            if "Commercial" in p.text:
                p.text = "Commercial / Technical / Financial viability: Highly viable SaaS model. Significant potential in 6 crore+ Indian MSMEs and global developing economies."
            elif "Import Substitution" in p.text:
                p.text = "Make in India / Atmanirbhar Bharat: Empowers local manufacturing to become globally competitive."
            elif "Energy saving" in p.text:
                p.text = "Process Optimization: AI-assisted process optimization supports efficiency and predictive maintenance."
            elif "Raw material" in p.text:
                p.text = "Resource efficiency: Enhances product quality by identifying recurring issues and reducing waste."

    # Slide 11: Impact on the project
    s11 = prs.slides[10]
    body11 = get_body_shape(s11)
    if body11:
        for p in body11.text_frame.paragraphs:
            if "Climate Mitigation" in p.text:
                p.text = "Environmental / Operational Impact: Enables digital transformation without heavy infrastructure, promoting efficient resource usage."
            elif "Intellectual property" in p.text:
                p.text = "Status of Intellectual property: FactoryDNA builds a proprietary centralized knowledge repository for the factory's exclusive historical knowledge."

    prs.save(output_path)
    print(f"Presentation successfully updated and saved to {output_path}")

if __name__ == "__main__":
    populate_ppt_carefully(sys.argv[1], sys.argv[2])
