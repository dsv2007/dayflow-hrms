import sys
from pptx import Presentation

def inspect_budget(file_path):
    prs = Presentation(file_path)
    slide = prs.slides[11]
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            print(f"Table: {len(table.rows)} rows, {len(table.columns)} cols")
            for r_idx, row in enumerate(table.rows):
                row_data = []
                for c_idx, cell in enumerate(row.cells):
                    # just print a snippet of text to identify the cell
                    text = cell.text.replace('\n', ' ')[:30]
                    row_data.append(f"({r_idx},{c_idx}): '{text}'")
                print(" | ".join(row_data))

if __name__ == "__main__":
    inspect_budget(sys.argv[1])
