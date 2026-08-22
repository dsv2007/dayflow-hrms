import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def add_flowchart(template_path, output_path):
    prs = Presentation(template_path)
    
    # Slide 4 (index 3) is the Flow Chart slide
    slide = prs.slides[3]
    
    # First, let's remove the placeholder text box to make room for the flowchart
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape != slide.shapes.title:
            # We can't directly delete, but we can move it off-screen or clear it
            shape.text = ""
            # Actually, to be safe, just clear text. We will draw over it.

    # Flowchart blocks
    # 1. Inputs (Top)
    left1 = Inches(1.5)
    top1 = Inches(2.0)
    width1 = Inches(3.0)
    height1 = Inches(1.0)
    
    left2 = Inches(5.5)
    top2 = Inches(2.0)
    width2 = Inches(3.0)
    height2 = Inches(1.0)
    
    # 2. AI Engine (Middle)
    left3 = Inches(3.5)
    top3 = Inches(3.5)
    width3 = Inches(3.0)
    height3 = Inches(1.0)
    
    # 3. Knowledge Base (Bottom)
    left4 = Inches(3.5)
    top4 = Inches(5.0)
    width4 = Inches(3.0)
    height4 = Inches(1.0)
    
    # 4. Output (Bottom-most)
    left5 = Inches(2.5)
    top5 = Inches(6.5)
    width5 = Inches(5.0)
    height5 = Inches(0.8)

    # Function to create a styled block
    def create_block(slide, left, top, width, height, text, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, rgb=(0, 102, 204)):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1.5)
        
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        return shape

    def create_arrow(slide, start_x, start_y, end_x, end_y):
        shape = slide.shapes.add_connector(
            MSO_SHAPE.DOWN_ARROW, start_x, start_y, end_x, end_y
        )
        shape.line.color.rgb = RGBColor(100, 100, 100)
        shape.line.width = Pt(3)

    # Create Blocks
    b1 = create_block(slide, left1, top1, width1, height1, "Unstructured Data\n(Logs, SOPs, Images, Voice)", rgb=(20, 140, 100))
    b2 = create_block(slide, left2, top2, width2, height2, "Structured Data\n(IoT, ERP, MES Sensors)", rgb=(20, 140, 100))
    
    b3 = create_block(slide, left3, top3, width3, height3, "AI Processing Engine\n(RAG, Edge AI, Knowledge Graph)", rgb=(200, 80, 50))
    
    b4 = create_block(slide, left4, top4, width4, height4, "Centralized Knowledge Base\n(Searchable & Structured)", rgb=(50, 100, 200))
    
    b5 = create_block(slide, left5, top5, width5, height5, "AI Manufacturing Copilot\n(Contextual Guidance & Optimization)", shape_type=MSO_SHAPE.HEXAGON, rgb=(120, 50, 180))

    # We will use simple DOWN_ARROW shapes instead of connectors to make it easier and more attractive
    def add_down_arrow(slide, left, top, width=Inches(0.4), height=Inches(0.4)):
        shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        shape.line.color.rgb = RGBColor(100, 100, 100)

    # Arrow from b1 to b3
    add_down_arrow(slide, Inches(3.0), Inches(3.05))
    # Arrow from b2 to b3
    add_down_arrow(slide, Inches(6.5), Inches(3.05))
    # Arrow from b3 to b4
    add_down_arrow(slide, Inches(4.8), Inches(4.55))
    # Arrow from b4 to b5
    add_down_arrow(slide, Inches(4.8), Inches(6.05))

    prs.save(output_path)
    print(f"Flowchart successfully generated and saved to {output_path}")

if __name__ == "__main__":
    add_flowchart(sys.argv[1], sys.argv[2])
