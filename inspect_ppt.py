import sys
from pptx import Presentation

def inspect_ppt(file_path):
    prs = Presentation(file_path)
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape != slide.shapes.title:
                    print(f"Text content: {shape.text[:100]}...")

if __name__ == "__main__":
    inspect_ppt(sys.argv[1])
